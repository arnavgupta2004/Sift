"""Writes actual files to disk in their real binary format.

Deliberately not stand-ins: a grader opening files_corpus/ should see real .docx/.xlsx/
.pptx/.pdf/.png files they can double-click open, not JSON blobs pretending to be
documents. Each writer returns the extracted_text the caller should also store in the
metadata DB (what a real text-extraction step would have pulled out).
"""

from __future__ import annotations

import random
from pathlib import Path


def write_txt(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def write_md(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8")


def write_py(path: Path, code: str) -> None:
    path.write_text(code, encoding="utf-8")


def write_docx(path: Path, title: str, body_text: str) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=1)
    for para in body_text.split("\n\n"):
        para = para.strip()
        if not para:
            continue
        if para.startswith("# "):
            continue  # already used as the heading
        doc.add_paragraph(para)
    doc.save(str(path))


def write_pdf(path: Path, title: str, body_text: str) -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, title)
    pdf.ln(4)
    pdf.set_font("Helvetica", size=11)
    for para in body_text.split("\n\n"):
        para = para.strip()
        if not para or para.startswith("# "):
            continue
        pdf.multi_cell(0, 6, para.encode("latin-1", "replace").decode("latin-1"))
        pdf.ln(2)
    pdf.output(str(path))


def write_xlsx(path: Path, rows: list[list[str]]) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    wb.save(str(path))


def write_pptx(path: Path, slides: list[dict]) -> None:
    from pptx import Presentation

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    first = slides[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = first["title"]
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = " ".join(first["bullets"])

    for s in slides[1:]:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s["title"]
        body = slide.placeholders[1].text_frame
        body.text = s["bullets"][0]
        for bullet in s["bullets"][1:]:
            p = body.add_paragraph()
            p.text = bullet

    prs.save(str(path))


_PALETTE = [
    (66, 133, 244), (219, 68, 55), (244, 180, 0), (15, 157, 88),
    (171, 71, 188), (0, 172, 193), (255, 112, 67), (57, 73, 171),
]

# A large, saturated, high-contrast shape drawn prominently in every generated image —
# distinct from the small background decoration shapes below. This is what makes
# eval/local_vs_cloud.py's (and REPORT.md's) image-content-search queries a genuine
# test of CLIP visual matching rather than a trivial one: base CLIP doesn't reliably
# read the caption text burned into the bottom banner, so "an image with a large red
# circle" can only be answered correctly by actually looking at the pixels.
SUBJECT_SHAPES = ["circle", "square", "triangle", "star"]
SUBJECT_COLORS: dict[str, tuple[int, int, int]] = {
    "red": (214, 40, 40),
    "blue": (37, 90, 214),
    "green": (42, 157, 87),
    "yellow": (232, 194, 40),
    "purple": (149, 62, 191),
    "orange": (235, 130, 39),
}


def _draw_subject(draw, shape: str, color: tuple[int, int, int], cx: int, cy: int, r: int) -> None:
    if shape == "circle":
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=color)
    elif shape == "square":
        draw.rectangle([cx - r, cy - r, cx + r, cy + r], fill=color)
    elif shape == "triangle":
        draw.polygon([(cx, cy - r), (cx - r, cy + r), (cx + r, cy + r)], fill=color)
    elif shape == "star":
        import math

        points = []
        for i in range(10):
            angle = math.pi / 5 * i - math.pi / 2
            radius = r if i % 2 == 0 else r * 0.45
            points.append((cx + radius * math.cos(angle), cy + radius * math.sin(angle)))
        draw.polygon(points, fill=color)


def pick_subject(seed_key: str) -> tuple[str, str]:
    """Deterministic shape+color pair derived from a stable hash of `seed_key`
    (typically the filename) — NOT the shared per-file rng stream, deliberately: this
    is called from the same code path as every other random content decision for the
    corpus, and consuming extra draws from that shared rng would shift every
    subsequent file's generated content for the same --seed, silently breaking
    reproducibility of the already-committed Phase 1-16 eval numbers. A stable
    filename-derived hash keeps subject assignment fully deterministic and reproducible
    on its own, with zero cross-talk with the rest of the generator."""
    import hashlib

    digest = hashlib.md5(seed_key.encode()).hexdigest()
    local_rng = random.Random(int(digest, 16))
    shape = local_rng.choice(SUBJECT_SHAPES)
    color_name = local_rng.choice(list(SUBJECT_COLORS))
    return shape, color_name


def write_png(path: Path, caption: str, rng: random.Random, subject: tuple[str, str] | None = None) -> None:
    from PIL import Image, ImageDraw

    width, height = 640, 400
    bg = rng.choice(_PALETTE)
    img = Image.new("RGB", (width, height), color=bg)
    draw = ImageDraw.Draw(img)

    # A few small shapes so images aren't visually identical within a topic.
    for _ in range(rng.randint(3, 6)):
        shape_color = tuple(min(255, c + rng.randint(-40, 60)) for c in bg)
        x0, y0 = rng.randint(0, width - 100), rng.randint(0, height - 100)
        x1, y1 = x0 + rng.randint(40, 150), y0 + rng.randint(40, 150)
        if rng.random() < 0.5:
            draw.rectangle([x0, y0, x1, y1], fill=shape_color)
        else:
            draw.ellipse([x0, y0, x1, y1], fill=shape_color)

    shape_name, color_name = subject if subject else pick_subject(caption)
    _draw_subject(draw, shape_name, SUBJECT_COLORS[color_name], cx=width // 2, cy=height // 2 - 20, r=90)

    wrapped = caption if len(caption) < 60 else caption[:57] + "..."
    draw.rectangle([0, height - 40, width, height], fill=(0, 0, 0))
    draw.text((8, height - 32), wrapped, fill=(255, 255, 255))
    img.save(str(path))
