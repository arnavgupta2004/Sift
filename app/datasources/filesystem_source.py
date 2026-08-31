"""Crawls a real local directory tree into RawFile records — real, messy, unlabeled
data (no topic_cluster ground truth; files get "uncategorized", and
app.personalization.profile_builder's discovered KMeans clusters still work fine since
they're computed from embeddings, not this label).

Text extraction is real, not a stub: native text formats are read directly, and
.pdf/.docx/.xlsx/.pptx go through the same libraries used elsewhere in this repo
(pypdf, python-docx, openpyxl, python-pptx). Anything else (images, archives,
binaries) still gets indexed with real filename/metadata — just no semantic content,
so it's findable by filename/metadata search but not by keyword/semantic search.
"""

from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path

from app.datasources.base import DataSource, RawFile

TEXT_NATIVE_EXTENSIONS = {"txt", "md", "py", "json", "csv", "yaml", "yml", "js", "ts", "tsx", "jsx", "html", "css", "sh"}
MAX_CONTENT_CHARS = 20_000
DEFAULT_EXCLUDE_DIRS = {".git", "node_modules", "__pycache__", ".venv", "venv", "dist", ".pytest_cache", "data"}


class FilesystemDataSource(DataSource):
    def __init__(
        self,
        root: Path | str,
        max_files: int | None = None,
        exclude_dirs: set[str] | None = None,
    ):
        self.root = Path(root).expanduser().resolve()
        self.max_files = max_files
        self.exclude_dirs = exclude_dirs if exclude_dirs is not None else DEFAULT_EXCLUDE_DIRS

    def list_files(self) -> list[RawFile]:
        if not self.root.exists():
            raise FileNotFoundError(f"root does not exist: {self.root}")

        results: list[RawFile] = []
        for dirpath, dirnames, filenames in os.walk(self.root):
            dirnames[:] = [d for d in dirnames if d not in self.exclude_dirs and not d.startswith(".")]
            for name in filenames:
                if name.startswith("."):
                    continue
                full_path = Path(dirpath) / name
                try:
                    stat = full_path.stat()
                except OSError:
                    continue

                ext = full_path.suffix.lstrip(".").lower() or "unknown"
                rel_path = str(full_path.relative_to(self.root))
                text = _extract_text(full_path, ext)

                results.append(
                    RawFile(
                        filename=name,
                        path=rel_path,
                        file_type=ext,
                        size_bytes=stat.st_size,
                        created_at=datetime.fromtimestamp(stat.st_ctime),
                        modified_at=datetime.fromtimestamp(stat.st_mtime),
                        extracted_text=text,
                    )
                )
                if self.max_files and len(results) >= self.max_files:
                    return results
        return results


def _extract_text(path: Path, ext: str) -> str:
    try:
        if ext in TEXT_NATIVE_EXTENSIONS:
            return path.read_text(encoding="utf-8", errors="ignore")[:MAX_CONTENT_CHARS]
        if ext == "pdf":
            return _extract_pdf(path)
        if ext == "docx":
            return _extract_docx(path)
        if ext == "xlsx":
            return _extract_xlsx(path)
        if ext == "pptx":
            return _extract_pptx(path)
    except Exception:
        return ""
    return ""


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages[:20]:  # cap pages for very long PDFs
        parts.append(page.extract_text() or "")
    return "\n".join(parts)[:MAX_CONTENT_CHARS]


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)[:MAX_CONTENT_CHARS]


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            lines.append(", ".join(str(c) for c in row if c is not None))
            if len(lines) > 500:
                break
    return "\n".join(lines)[:MAX_CONTENT_CHARS]


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)
    return "\n".join(lines)[:MAX_CONTENT_CHARS]
